# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==================== 配色 ====================
BG = RGBColor(0x1A, 0x1A, 0x2E)
GREEN = RGBColor(0x00, 0xE6, 0x76)
BLUE = RGBColor(0x00, 0xB4, 0xD8)
ORANGE = RGBColor(0xFF, 0x9F, 0x1C)
RED = RGBColor(0xFF, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xCC)
CARD = RGBColor(0x25, 0x25, 0x3D)
LINE = RGBColor(0x3A, 0x3A, 0x55)
YELLOW = RGBColor(0xFF, 0xE0, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width


def bg_fill(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s


def card(slide, l, t, w, h, c=CARD, r=0.03):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if r:
        s.adjustments[0] = r
    return s


def txt(slide, text, l, t, w, h, sz=18, c=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.alignment = align
    return tb


def lines(slide, text_lines, l, t, w, h, sz=15, c=GRAY, sp=1.6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(sz * (sp - 1))
    return tb


def circle(slide, l, t, sz, c, text, fsz=20):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    p = s.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(fsz); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    return s


def slide_header(slide, title, accent_color):
    rect(slide, Inches(0), Inches(0), SW, Inches(0.06), accent_color)
    txt(slide, title, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
        sz=32, bold=True)
    rect(slide, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.05), accent_color)


# ==================== P1: 封面 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
rect(s, Inches(0), Inches(0), SW, Inches(0.06), GREEN)
txt(s, "AI 动作姿态纠错系统", Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
    sz=48, bold=True, align=PP_ALIGN.CENTER)
txt(s, "输入一段老师视频和一段学生视频", Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
    sz=22, c=GRAY, align=PP_ALIGN.CENTER)
txt(s, "自动输出一段左右对比的纠错视频", Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
    sz=22, c=GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(5), Inches(5.4), Inches(3.333), Inches(0.02), LINE)


# ==================== P2: 为什么要做这个 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "项目背景", ORANGE)

card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(1.1), Inches(2.1), Inches(0.08), Inches(0.5), RED)
txt(s, "传统教学的痛点", Inches(1.35), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=RED, bold=True)
lines(s, [
    "  老师示范动作, 学生模仿练习",
    "",
    "  但学生自己看不到自己的动作",
    "",
    "  没有老师在旁边, 就没人纠正",
    "",
    "  错误的动作反复练习, 越练越错",
], Inches(1.1), Inches(2.9), Inches(4.8), Inches(3.5), sz=17, sp=1.7)

card(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(7.3), Inches(2.1), Inches(0.08), Inches(0.5), GREEN)
txt(s, "我们的思路", Inches(7.55), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=GREEN, bold=True)
lines(s, [
    "  用 AI 自动比对老师和学生的动作",
    "",
    "  像一面\"智能镜子\":",
    "  告诉你哪里做错了, 该怎么改",
    "",
    "  不需要老师实时在场",
    "",
    "  学生自己就能练习和纠错",
], Inches(7.3), Inches(2.9), Inches(4.8), Inches(3.5), sz=17, sp=1.7)


# ==================== P3: 系统做了什么 (效果展示) ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "系统效果", GREEN)

# 左边：输入
card(s, Inches(0.8), Inches(1.8), Inches(3.5), Inches(5.0))
rect(s, Inches(1.1), Inches(2.1), Inches(0.08), Inches(0.5), BLUE)
txt(s, "输入", Inches(1.35), Inches(2.1), Inches(2.8), Inches(0.5),
    sz=22, c=BLUE, bold=True)
lines(s, [
    "  老师示范视频",
    "",
    "  学生练习视频",
    "",
    "  (任意角度, 任意速度)",
], Inches(1.1), Inches(3.0), Inches(2.8), Inches(3.0), sz=16, sp=1.8)

# 箭头
txt(s, ">>>", Inches(4.5), Inches(3.8), Inches(1.2), Inches(0.8),
    sz=36, c=GREEN, bold=True, align=PP_ALIGN.CENTER)

# 右边：输出
card(s, Inches(5.8), Inches(1.8), Inches(6.7), Inches(5.0))
rect(s, Inches(6.1), Inches(2.1), Inches(0.08), Inches(0.5), GREEN)
txt(s, "输出: 纠错视频", Inches(6.35), Inches(2.1), Inches(5.8), Inches(0.5),
    sz=22, c=GREEN, bold=True)
lines(s, [
    "  左右拼接的对比视频:",
    "",
    "  左半边: 老师视频 + 绿色骨架 (标准参考)",
    "",
    "  右半边: 学生视频 + 三层叠加:",
    "    - 老师残影 (半透明绿色, 作为参考)",
    "    - 学生骨架 (正确=橙色, 错误=红色)",
    "    - 纠错箭头 (指出正确方向)",
], Inches(6.1), Inches(3.0), Inches(5.8), Inches(3.5), sz=15, sp=1.6)


# ==================== P4: 整体流程一览 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "处理流程一览", BLUE)

steps = [
    ("1", "提取骨架", "MediaPipe\n33个关节点", GREEN),
    ("2", "平滑降噪", "滑动平均\n消除抖动", YELLOW),
    ("3", "空间归一化", "消除体型差异\n统一坐标尺度", ORANGE),
    ("4", "时间对齐", "DTW算法\n解决速度差异", BLUE),
    ("5", "误差检测", "逐关节比对\n标记错误", RED),
    ("6", "渲染输出", "骨架+箭头\n合成视频", GREEN),
]

bw = Inches(1.7); bh = Inches(2.8); gap = Inches(0.25)
tw = len(steps) * bw + (len(steps)-1) * gap
sx = int((SW - tw) / 2); sy = Inches(2.2)

for i, (n, title, desc, c) in enumerate(steps):
    x = sx + int(i * (bw + gap))
    card(s, x, sy, bw, bh, r=0.05)
    rect(s, x, sy, bw, Inches(0.06), c)
    circle(s, x + Inches(0.6), sy + Inches(0.3), Inches(0.5), c, n, 18)
    txt(s, title, x+Inches(0.1), sy+Inches(1.0), Inches(1.5), Inches(0.5),
        sz=16, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc, x+Inches(0.1), sy+Inches(1.6), Inches(1.5), Inches(1.0),
        sz=12, c=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        txt(s, ">", x+bw+Inches(0.02), sy+Inches(1.0), Inches(0.22), Inches(0.5),
            sz=22, c=GREEN, bold=True, align=PP_ALIGN.CENTER)

txt(s, "下面重点讲三个核心技术难点",
    Inches(2), Inches(5.8), Inches(9.333), Inches(0.5),
    sz=18, c=ORANGE, bold=True, align=PP_ALIGN.CENTER)


# ==================== P5: 难点一 - 体型不同怎么办 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "难点一: 体型不同怎么比?", ORANGE)

# 问题
card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(1.1), Inches(2.1), Inches(0.08), Inches(0.5), RED)
txt(s, "问题", Inches(1.35), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=RED, bold=True)
lines(s, [
    "  老师和学生身高体型不同",
    "",
    "  同一个\"手臂平举\"动作",
    "  手腕的绝对坐标完全不一样",
    "",
    "  直接比对坐标 -> 全是误报",
], Inches(1.1), Inches(2.9), Inches(4.8), Inches(3.5), sz=17, sp=1.8)

# 方案
card(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(7.3), Inches(2.1), Inches(0.08), Inches(0.5), GREEN)
txt(s, "方案: 用自己的身体当尺子", Inches(7.55), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=GREEN, bold=True)
lines(s, [
    "  1. 算出骨盆中心点 (作为原点)",
    "",
    "  2. 算出躯干长度 (作为单位长度)",
    "",
    "  3. 所有关节坐标 = (原坐标 - 原点) / 躯干长度",
    "",
    "  归一化后:",
    "  不管高矮胖瘦, 做同一个动作",
    "  得到的数值基本一致",
], Inches(7.3), Inches(2.9), Inches(4.8), Inches(3.8), sz=16, sp=1.6)


# ==================== P6: 难点二 - 速度不同怎么对齐 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "难点二: 速度不同怎么对齐?", BLUE)

card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(1.1), Inches(2.1), Inches(0.08), Inches(0.5), RED)
txt(s, "问题", Inches(1.35), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=RED, bold=True)
lines(s, [
    "  老师和学生做同一套动作",
    "  但节奏不可能完全一致",
    "",
    "  学生可能某个动作做得更慢",
    "",
    "  逐帧比对 -> 帧错位 -> 全是误报",
], Inches(1.1), Inches(2.9), Inches(4.8), Inches(3.5), sz=17, sp=1.8)

card(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(7.3), Inches(2.1), Inches(0.08), Inches(0.5), GREEN)
txt(s, "方案: DTW 动态时间规整", Inches(7.55), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=GREEN, bold=True)
lines(s, [
    "  一种经典的序列对齐算法",
    "",
    "  自动找到两个序列的",
    "  \"最优帧对应关系\"",
    "",
    "  允许一对多, 多对一的映射",
    "",
    "  效果: 不管快做慢做",
    "  同一动作的帧都能正确匹配",
], Inches(7.3), Inches(2.9), Inches(4.8), Inches(3.8), sz=16, sp=1.6)


# ==================== P7: 难点三 - 怎么判断做错了 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "难点三: 怎么判断做错了?", RED)

card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(1.1), Inches(2.1), Inches(0.08), Inches(0.5), ORANGE)
txt(s, "误差检测", Inches(1.35), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=ORANGE, bold=True)
lines(s, [
    "  对每个关节, 计算学生和老师的",
    "  归一化坐标的欧氏距离",
    "",
    "  距离 > 阈值 (躯干长度的12%)",
    "  -> 标记为错误",
    "",
    "  只检测关键末端关节:",
    "  手腕/手肘/膝盖/脚踝",
], Inches(1.1), Inches(2.9), Inches(4.8), Inches(3.5), sz=17, sp=1.7)

card(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
rect(s, Inches(7.3), Inches(2.1), Inches(0.08), Inches(0.5), YELLOW)
txt(s, "防闪烁机制", Inches(7.55), Inches(2.1), Inches(4.5), Inches(0.5),
    sz=22, c=YELLOW, bold=True)
lines(s, [
    "  问题: 某一帧刚好在阈值边缘",
    "  -> 红一下灭一下, 不停闪烁",
    "",
    "  解决: 状态保持 10 帧",
    "  一旦触发错误, 至少持续 10 帧",
    "  即使中间误差变小也不消失",
    "",
    "  效果: 错误标注稳定, 不闪烁",
], Inches(7.3), Inches(2.9), Inches(4.8), Inches(3.8), sz=16, sp=1.6)


# ==================== P8: 视觉渲染 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "视觉渲染: 学生画面上的三层叠加", GREEN)

# 三层
layer_data = [
    ("第 1 层: 老师残影", "对齐后的老师骨架, 半透明绿色, 作为参考标准", GREEN),
    ("第 2 层: 学生骨架", "正确部位 = 橙色, 错误部位 = 红色粗线高亮", ORANGE),
    ("第 3 层: 纠错箭头", "从学生错误位置 -> 老师正确位置 (红色箭头)", RED),
]

y_start = Inches(1.8)
for i, (title, desc, c) in enumerate(layer_data):
    y = y_start + Inches(i * 1.8)
    card(s, Inches(0.8), y, Inches(11.7), Inches(1.5))
    rect(s, Inches(1.1), y + Inches(0.3), Inches(0.08), Inches(0.5), c)
    txt(s, title, Inches(1.4), y + Inches(0.25), Inches(4), Inches(0.5),
        sz=20, c=c, bold=True)
    txt(s, desc, Inches(1.4), y + Inches(0.8), Inches(10), Inches(0.5),
        sz=16, c=GRAY)

# 底部补充
txt(s, "箭头长度 < 15px 时不画, 避免微小偏差导致画面抖动",
    Inches(2), Inches(7.0), Inches(9.333), Inches(0.4),
    sz=14, c=GRAY, align=PP_ALIGN.CENTER)


# ==================== P9: 总结 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
slide_header(s, "总结", GREEN)

card(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(5.0))

lines(s, [
    "  本项目实现了一个全自动的动作纠错系统:",
    "",
    "  - 输入: 老师视频 + 学生视频",
    "  - 输出: 带骨架标注和纠错箭头的对比视频",
    "",
    "  解决了三个核心技术难点:",
    "",
    "    1. 体型差异  ->  空间归一化 (用自己的身体当尺子)",
    "    2. 速度差异  ->  DTW 时间对齐 (自动匹配同动作帧)",
    "    3. 稳定显示  ->  状态保持 + 三层渲染 (不闪不抖)",
    "",
    "  适用场景: 健身/舞蹈/体育/康复训练 等",
], Inches(2.0), Inches(2.2), Inches(9), Inches(4.2), sz=18, sp=1.5)


# ==================== P10: 谢谢 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_fill(s)
rect(s, Inches(0), Inches(0), SW, Inches(0.06), GREEN)
txt(s, "Thank You", Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
    sz=56, bold=True, align=PP_ALIGN.CENTER)
txt(s, "谢谢观看", Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
    sz=28, c=GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(5), Inches(5.2), Inches(3.333), Inches(0.02), LINE)


# ==================== 保存 ====================
out = r"e:\leetcode\last\AI_Pose_Correction_Intro_v5.pptx"
prs.save(out)
print(f"PPT saved: {out}")
